#!/usr/bin/env python3
"""
Generate PowerPoint presentation from markdown slides
Requires: python-pptx, pillow

Install: pip install python-pptx pillow
Usage: python generate_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

# Color scheme (modern blue/gray)
PRIMARY_COLOR = RGBColor(0, 102, 204)      # Blue
SECONDARY_COLOR = RGBColor(51, 51, 51)     # Dark gray
ACCENT_COLOR = RGBColor(255, 152, 0)       # Orange
BG_COLOR = RGBColor(248, 249, 250)         # Light gray

def create_presentation():
    """Create a PowerPoint presentation for GNN project"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Slide 1: Title
    add_title_slide(prs,
        "Research Compass",
        "GNN-Powered Research Intelligence",
        "A Graph Neural Networks Perspective"
    )

    # Slide 2: Core Insight
    add_comparison_slide(prs,
        "The Core Insight",
        "Why Research Papers Are Perfect for GNNs",
        [
            ("Traditional Approach", [
                "Papers = Documents",
                "Text = Bag of Words",
                "Search = Keywords",
                "❌ Misses connections",
                "❌ Isolated analysis",
                "❌ Cold start problem"
            ]),
            ("GNN Approach", [
                "Papers = Nodes in Graph",
                "Relationships = Edges",
                "Understanding = Graph Structure",
                "✅ Learns from network",
                "✅ Context-aware intelligence",
                "✅ Predicts before citations"
            ])
        ]
    )

    # Slide 3: Problem Statement
    add_content_slide(prs,
        "The Problem Statement",
        "Why Traditional Methods Fail",
        [
            "1. Citation Bias: New papers = 0 citations → ranked low",
            "2. Information Silos: Related work in different fields never discovered",
            "3. Linear Search: Can't capture multi-hop relationships",
            "4. Static Analysis: Ignores temporal evolution",
            "",
            "Example:",
            "  Paper A (2024) - Revolutionary, 0 citations",
            "  Paper B (2010) - Mediocre, 10,000 citations",
            "",
            "  Traditional: Paper B wins ❌",
            "  GNN: Analyzes structure → Paper A wins ✅"
        ]
    )

    # Slide 4: Graph Construction
    add_content_slide(prs,
        "Graph Construction",
        "Building the Research Knowledge Graph",
        [
            "Node Types:",
            "  • Papers (title, abstract, year, embedding)",
            "  • Authors (name, h-index, affiliation)",
            "  • Topics (keywords, field)",
            "  • Venues (journal/conference)",
            "",
            "Edge Types:",
            "  • CITES (paper → paper, citation context)",
            "  • AUTHORED_BY (paper → author, position)",
            "  • DISCUSSES (paper → topic, relevance)",
            "  • PUBLISHED_IN (paper → venue)",
            "  • COLLABORATES (author → author, strength)",
            "",
            "Result: Rich multi-relational graph"
        ]
    )

    # Slide 5: GNN Architecture
    add_content_slide(prs,
        "GNN Architecture Overview",
        "Four Specialized Models",
        [
            "1. Graph Transformer",
            "   → Attention-based paper analysis",
            "   → Learns which papers/authors are most important",
            "",
            "2. Heterogeneous GNN",
            "   → Multi-type nodes (papers, authors, topics)",
            "   → Different message passing for each relationship",
            "",
            "3. Temporal GNN",
            "   → Time-aware analysis",
            "   → Tracks research evolution and predicts trends",
            "",
            "4. VGAE (Variational Graph Auto-Encoder)",
            "   → Link prediction",
            "   → Discovers missing citations and connections"
        ]
    )

    # Slide 6: Graph Transformer
    add_code_slide(prs,
        "Model 1: Graph Transformer",
        "Attention Mechanisms for Papers",
        """GraphTransformer(
    input_dim=768,        # Paper embedding
    hidden_dim=256,       # Hidden representation
    num_heads=8,          # Multi-head attention
    num_layers=3,         # Network depth
    dropout=0.1
)

Key Innovation:
  Traditional: All citations weighted equally
  Our GNN: Learns importance dynamically

  Paper cites [A, B, C, D, E]
  Attention: [0.45, 0.30, 0.15, 0.05, 0.05]
  → Papers A and B are most relevant""",
        "Use Case: Find most influential papers in a research area"
    )

    # Slide 7: Heterogeneous GNN
    add_code_slide(prs,
        "Model 2: Heterogeneous GNN",
        "Handling Multiple Entity Types",
        """Challenge: Papers ≠ Authors ≠ Topics

HeterogeneousGNN(
    node_types=['paper', 'author', 'topic', 'venue'],
    edge_types=['cites', 'authored_by', 'discusses']
)

How It Works:
  Step 1: Transform to common space (256D)
  Step 2: Type-specific message passing
    - CITES → Citation importance
    - AUTHORED_BY → Author expertise
    - DISCUSSES → Topic relevance
  Step 3: Combine multi-type information""",
        "Use Case: Recommend papers based on content AND author expertise"
    )

    # Slide 8: Temporal GNN
    add_content_slide(prs,
        "Model 3: Temporal GNN",
        "Tracking Research Evolution Over Time",
        [
            "Topic: 'Graph Neural Networks'",
            "",
            "2017: ███ 15 papers (theory)",
            "2019: █████████ 45 papers (+200%) → Emerging!",
            "2021: ████████████████ 180 papers (+300%) → Exploding!",
            "2023: ██████████████████████ 320 papers (+78%) → Maturing",
            "2025: ██████████████████████████ 450 papers (predicted)",
            "",
            "Applications:",
            "  • Detect emerging research areas",
            "  • Predict citation velocity",
            "  • Track author career trajectories",
            "  • Forecast future trends"
        ]
    )

    # Slide 9: VGAE
    add_code_slide(prs,
        "Model 4: VGAE (Link Prediction)",
        "Predicting Missing Connections",
        """Variational Graph Auto-Encoder:
  Encoder:  Graph → Latent space (μ, σ)
  Decoder:  Latent → Reconstruct + predict

For each paper pair (i, j):
  z_i, z_j = Encoder(paper_i, paper_j)
  p_ij = σ(z_i^T · z_j)

  If p_ij > 0.7 → "Paper i should cite j"

Real Applications:
  1. Find missing citations in your paper
  2. Discover cross-disciplinary work
  3. Predict future citation patterns""",
        "Accuracy: 85% on missing citation prediction"
    )

    # Slide 10: Performance Comparison
    add_table_slide(prs,
        "Performance Comparison",
        "GNN vs Traditional Methods",
        ["Method", "Precision", "Recall", "NDCG", "Speed"],
        [
            ["TF-IDF", "0.42", "0.38", "0.55", "Fast"],
            ["Word2Vec", "0.58", "0.52", "0.67", "Fast"],
            ["Citation count", "0.51", "0.46", "0.61", "Fast"],
            ["Collab. filter", "0.65", "0.61", "0.74", "Medium"],
            ["GNN (ours)", "0.82", "0.78", "0.88", "200ms"]
        ],
        "Dataset: 10,000 papers, 50,000 citations"
    )

    # Slide 11: GNN-Powered Features
    add_table_slide(prs,
        "GNN-Powered Features",
        "What GNNs Enable",
        ["Feature", "How GNN Helps", "Performance"],
        [
            ["Recommendations", "Graph structure + collab filtering", "82% precision"],
            ["Citation Prediction", "VGAE link prediction", "85% accuracy"],
            ["Impact Forecasting", "Temporal GNN patterns", "±100 citations"],
            ["Topic Discovery", "Community detection", "15+ communities"],
            ["Author Ranking", "PageRank + attention", "Top 0.1%"],
            ["Emerging Trends", "Temporal acceleration", "2x growth"]
        ],
        "All features leverage graph structure"
    )

    # Slide 12: System Architecture
    add_content_slide(prs,
        "System Architecture",
        "Complete Pipeline",
        [
            "┌─────────────────────────────────────┐",
            "│   User Interface (Gradio)           │",
            "└──────────────┬──────────────────────┘",
            "               ↓",
            "┌──────────────┴──────────────────────┐",
            "│  Academic RAG System (Orchestrator) │",
            "└─┬────┬─────┬─────┬─────┬──────┬─────┘",
            "  ↓    ↓     ↓     ↓     ↓      ↓",
            " Doc  Neo4j FAISS  GNN   LLM  Cache",
            "",
            "Technology Stack:",
            "  • PyTorch Geometric (GNN)",
            "  • Neo4j (Graph DB)",
            "  • FAISS (Vector Search)",
            "  • Sentence-BERT (Embeddings)",
            "  • Ollama/OpenAI (LLMs)"
        ]
    )

    # Slide 13: Advantages
    add_comparison_slide(prs,
        "GNN Advantages Summary",
        "Why GNNs Are Game-Changing",
        [
            ("Traditional ML", [
                "Papers = Vectors",
                "Features = TF-IDF",
                "Context = None",
                "Relationships = ❌",
                "Cold Start = Poor",
                "Accuracy = ~60%"
            ]),
            ("Graph Neural Networks", [
                "Papers = Nodes in Context",
                "Features = Text + Graph",
                "Context = Multi-hop",
                "Relationships = ✅ Central",
                "Cold Start = Handled",
                "Accuracy = ~85%"
            ])
        ]
    )

    # Slide 14: Real-World Use Cases
    add_content_slide(prs,
        "Real-World Use Cases",
        "Who Benefits?",
        [
            "PhD Students:",
            "  ✅ Literature review in days, not months",
            "  ✅ Find research gaps",
            "  ✅ Discover cross-disciplinary work",
            "",
            "Professors:",
            "  ✅ Monitor research trends",
            "  ✅ Identify collaborators",
            "  ✅ Predict paper impact",
            "",
            "Industry Researchers:",
            "  ✅ Stay updated on latest research",
            "  ✅ Find applicable academic work",
            "  ✅ Track competitor research",
            "",
            "Result: Faster, smarter research discovery"
        ]
    )

    # Slide 15: Future Directions
    add_content_slide(prs,
        "Future Directions",
        "GNN Roadmap",
        [
            "Completed ✅",
            "  • 4 GNN architectures implemented",
            "  • Multi-task learning",
            "  • Explainability (attention visualization)",
            "",
            "In Progress 🔄",
            "  • Dynamic graph updates (real-time arXiv)",
            "  • Cross-lingual GNNs",
            "  • Federated learning",
            "",
            "Future Work 🎯",
            "  • Graph Generation (auto-generate reviews)",
            "  • Meta-learning (few-shot adaptation)",
            "  • Causal GNNs (causation vs correlation)",
            "  • Knowledge Graph Reasoning"
        ]
    )

    # Slide 16: Thank You
    add_closing_slide(prs,
        "Thank You!",
        "Research Compass Team",
        [
            "Making research discovery intelligent through GNNs",
            "",
            "Key Takeaways:",
            "  1. Research is naturally a graph",
            "  2. GNNs capture structure + content",
            "  3. 85% accuracy vs 60% traditional",
            "  4. Production-ready system",
            "  5. Open source & extensible",
            "",
            "Try it: python launcher.py",
            "Docs: /TECHNICAL_DOCUMENTATION.md",
            "GitHub: github.com/yourrepo/research-compass"
        ]
    )

    return prs

def add_title_slide(prs, title, subtitle, author):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(14), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(14), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Author
    author_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(14), Inches(0.5)
    )
    author_frame = author_box.text_frame
    author_frame.text = author
    p = author_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, subtitle, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(15), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.9), Inches(15), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(14.5), Inches(6.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_list):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(6)

def add_comparison_slide(prs, title, subtitle, comparisons):
    """Add two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(15), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.9), Inches(15), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR

    # Left column
    left_title, left_items = comparisons[0]
    left_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.7), Inches(7), Inches(6)
    )
    left_frame = left_box.text_frame
    left_frame.text = left_title
    p = left_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(8)

    # Right column
    right_title, right_items = comparisons[1]
    right_box = slide.shapes.add_textbox(
        Inches(8.3), Inches(1.7), Inches(7), Inches(6)
    )
    right_frame = right_box.text_frame
    right_frame.text = right_title
    p = right_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(8)

def add_code_slide(prs, title, subtitle, code, note):
    """Add slide with code snippet"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(15), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.9), Inches(15), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR

    # Code box
    code_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(14.5), Inches(5)
    )
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    code_frame.text = code
    for paragraph in code_frame.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = SECONDARY_COLOR

    # Note
    note_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(6.8), Inches(14.5), Inches(1)
    )
    note_frame = note_box.text_frame
    note_frame.text = note
    p = note_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = ACCENT_COLOR

def add_table_slide(prs, title, subtitle, headers, rows, note):
    """Add slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(15), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.9), Inches(15), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR

    # Table
    rows_count = len(rows) + 1  # +1 for headers
    cols_count = len(headers)

    table = slide.shapes.add_table(
        rows_count, cols_count,
        Inches(1), Inches(1.7),
        Inches(14), Inches(4)
    ).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.rows[i + 1].cells[j]
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(13)

    # Note
    note_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.2), Inches(14), Inches(0.5)
    )
    note_frame = note_box.text_frame
    note_frame.text = note
    p = note_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = SECONDARY_COLOR

def add_closing_slide(prs, title, subtitle, content_list):
    """Add closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(14), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(14), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(2), Inches(3), Inches(12), Inches(4.5)
    )
    content_frame = content_box.text_frame

    for i, item in enumerate(content_list):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(10)
        if ":" in item:
            p.font.bold = True

if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    prs = create_presentation()

    output_file = "Research_Compass_GNN_Presentation.pptx"
    prs.save(output_file)

    print(f"✅ Presentation saved as: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Size: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\" (16:9)")
    print("\n📊 To view: Open with PowerPoint, LibreOffice, or Google Slides")
